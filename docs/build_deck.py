"""Generate SENTINEL_Q_Deck.pptx. One-off build script (not part of the app)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BG = RGBColor(240, 248, 248)
TITLE_COLOR = RGBColor(10, 22, 40)
BODY_COLOR = RGBColor(10, 22, 40)
GOLD = RGBColor(245, 197, 24)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(10, 107, 122)
GREY = RGBColor(225, 225, 225)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    return slide


def set_font(run, size=18, bold=False, color=BODY_COLOR, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_title(slide, text, top=Inches(0.4), size=40):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=True, color=TITLE_COLOR)
    return box


def add_subtitle(slide, text, top=Inches(1.35), size=20):
    box = slide.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=False, color=TEAL)
    return box


def add_bullets(slide, bullets, top=Inches(2.2), left=Inches(0.8), width=Inches(11.7), size=18, height=Inches(4.5)):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = "•  " + b
        set_font(run, size=size, bold=False, color=BODY_COLOR)
    return box


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size=12, bold=False, color=TEAL)


def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    return shape


def add_placeholder_box(slide, left, top, width, height, label):
    shape = add_card(slide, left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_font(run, size=16, bold=True, color=TEAL)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


# ---------- SLIDE 1 ----------
s = add_slide()
box = slide_title = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.3))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "SENTINEL-Q"
set_font(run, size=54, bold=True, color=TITLE_COLOR)

box = s.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.9))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "The SOC sees the login. The fraud team sees the transfer. Nobody sees both."
set_font(run, size=22, bold=False, color=TEAL)

box = s.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.7))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Cyber-fraud correlation for Indian banks. Explainable. On-prem. Live right now."
set_font(run, size=18, bold=False, color=BODY_COLOR)

add_footer(s, "Finspark Hackathon 2026 · Problem Statement 2")

# ---------- SLIDE 2 ----------
s = add_slide()
add_title(s, "Banks Are Blind in the Middle")
add_subtitle(s, "Two teams, two dashboards, one attacker walking between them.")
add_bullets(s, [
    "The SIEM flags a suspicious login. The fraud engine flags a suspicious transfer. Neither knows the other fired.",
    "The attack lives in the gap: a compromised session that becomes a fraudulent transaction 40 minutes later.",
    "SENTINEL-Q correlates cybersecurity telemetry with transactional behaviour — one score, one screen, one story.",
])

# ---------- SLIDE 3 ----------
s = add_slide()
add_title(s, "What We Need to Run")
add_subtitle(s, "Almost nothing. That's the point.")
add_bullets(s, [
    "Pure Python. No cloud dependencies. Fully on-prem capable.",
    "Ingests two streams: security telemetry + transaction logs.",
    "Customer identifiers pseudonymized with SHA-256 before processing.",
    "Designed for compatibility with ISO 20022 (RTGS/NEFT) and NPCI UPI specifications.",
])

# ---------- SLIDE 4 ----------
s = add_slide()
add_title(s, "Built Lean. Built to Last.")
add_subtitle(s, "Five rules you can read. One model that learns. Nothing you can't audit.")
add_bullets(s, [
    "5 transparent correlation rules — every trigger condition is human-readable.",
    "IsolationForest anomaly weighting layered on top — catches what rules can't name.",
    "Output: a 0-100 risk score with a plain-English explanation of why it fired.",
    "No black box. An auditor can trace every score back to its cause.",
])

# ---------- SLIDE 5 ----------
s = add_slide()
add_title(s, "How a SOC Analyst Uses SENTINEL-Q")
add_subtitle(s, "From “something's off” to “here's exactly what happened.”")
add_bullets(s, [
    "An alert lands with a risk score and a plain-English reason — not a wall of logs.",
    "The analyst sees the cyber event and the transaction pattern side by side, already correlated.",
    "Five attack scenarios covered: account takeover, credential stuffing, dormant-account burst, HNDL exfiltration, impossible travel.",
])

# ---------- SLIDE 6 ----------
s = add_slide()
add_title(s, "Three Reasons Banks Can't Get This Elsewhere")
add_subtitle(s, "SIEM alone. FRM alone. Or both, finally talking.")

rows_data = [
    ["", "SIEM Alone", "FRM Alone", "SENTINEL-Q"],
    ["Sees the login", "YES", "NO", "YES"],
    ["Sees the transfer", "NO", "YES", "YES"],
    ["Correlates both", "NO", "NO", "YES"],
    ["Explains why", "NO", "NO", "YES"],
    ["Quantum exposure", "NO", "NO", "YES"],
]

table_top = Inches(2.3)
table_left = Inches(1.0)
table_width = Inches(11.3)
table_height = Inches(4.2)
gfx = s.shapes.add_table(len(rows_data), 4, table_left, table_top, table_width, table_height)
table = gfx.table

col_widths = [Inches(3.5), Inches(2.6), Inches(2.6), Inches(2.6)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = val
        bold = (r == 0)
        set_font(run, size=16, bold=bold, color=TITLE_COLOR if r == 0 else BODY_COLOR)

        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL
            run.font.color.rgb = WHITE
        else:
            if c == 3 and val == "YES":
                cell.fill.solid()
                cell.fill.fore_color.rgb = GOLD
            elif val == "NO":
                cell.fill.solid()
                cell.fill.fore_color.rgb = GREY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

# ---------- SLIDE 7 ----------
s = add_slide()
add_title(s, "The Code is Real. The Demo is Live.")
add_subtitle(s, "Not a mockup. Not a Figma. A running engine.")
add_bullets(s, [
    "Repo: github.com/Fablescript22/sentinel-q",
    "Live: sentinel-q-zjrt44jq29pa4uruuhbrkv.streamlit.app",
    "Deployed on Streamlit Community Cloud — open it on your phone right now.",
], top=Inches(2.2), height=Inches(2.2))
add_placeholder_box(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(2.2), "[SCREENSHOT PLACEHOLDER — insert dashboard screenshot here]")

# ---------- SLIDE 8 ----------
s = add_slide()
add_title(s, "[NUMBER] Fewer False Alarms. Every Day.")
add_subtitle(s, "Correlation kills noise. A login alert alone is not an incident.")
add_bullets(s, [
    "Rules fire only when cyber signal and transaction signal align — single-stream noise suppressed.",
    "[NUMBER] false positives suppressed in our test run — from engine terminal output.",
    "Fewer alerts, higher confidence per alert, less analyst fatigue.",
    "Assumption: 15 min triage per FP x Rs 800/hour analyst cost x 365 days = significant annual saving.",
])

# ---------- SLIDE 9 ----------
s = add_slide()
add_title(s, "Every Other Team Will Claim Quantum Detection. We Won't.")
add_subtitle(s, "You cannot detect a quantum computer. You can measure your exposure to one.")
add_bullets(s, [
    "Quantum Risk Monitor classifies crypto posture in 3 tiers: PQC-READY (ML-KEM FIPS 203), QUANTUM-VULNERABLE (ECDH), CRITICAL (RSA/legacy TLS).",
    "HNDL alerts fire ONLY when vulnerable crypto AND exfiltration-shaped behaviour occur together — never on crypto posture alone.",
    "We claim quantum exposure monitoring. Nothing more. That honesty is the feature.",
])

# ---------- SLIDE 10 ----------
s = add_slide()
add_title(s, "Three Clicks from Alert to Action")
add_subtitle(s, "Score. Explanation. Evidence.")
add_bullets(s, [
    "Click 1: the alert — 0-100 risk score, ranked by severity.",
    "Click 2: the explanation — plain English, which rules fired and why.",
    "Click 3: the evidence — correlated cyber event and transaction pattern behind the score.",
], top=Inches(2.2), height=Inches(2.2))
add_placeholder_box(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(2.2), "[SCREENSHOT PLACEHOLDER — insert alert detail screenshot here]")

# ---------- SLIDE 11 ----------
s = add_slide()
add_title(s, "One Engine. 50 Branches. No Redesign.")
add_subtitle(s, "Stateless by design. Scales by adding data, not code.")
add_bullets(s, [
    "Stateless batch engine — processes each branch independently, no shared state.",
    "Processes 13,000+ events per run on a single laptop — linear extrapolation to any volume.",
    "No database, no microservices, no infrastructure team required.",
    "Add a branch: point the engine at a new data feed. That is it.",
])

# ---------- SLIDE 12 ----------
s = add_slide()
add_title(s, "pip install. Done.")
add_subtitle(s, "If your bank runs Python, it runs SENTINEL-Q.")
add_bullets(s, [
    "Pure Python stack. No cloud services to procure, no vendor lock-in.",
    "One command to install: pip install -r requirements.txt",
    "Runs entirely on-prem — data never has to leave the bank.",
    "7 files. 6 libraries. No database. No Docker.",
])

# ---------- SLIDE 13 ----------
s = add_slide()
add_title(s, "Built for a Bank That Can't Afford a Breach")
add_subtitle(s, "Compliance is not a slide at the end. It is in the design.")
add_bullets(s, [
    "RBI Cyber Security Framework 2016 — SOC monitoring and incident detection aligned.",
    "DPDP Act 2023 — SHA-256 pseudonymization, data minimisation, purpose limitation.",
    "On-prem deployment keeps all data inside the bank perimeter.",
    "PQC readiness inventory via ML-KEM FIPS 203 / ML-DSA FIPS 204 tier classification.",
])

# ---------- SLIDE 14 ----------
s = add_slide()
add_title(s, "How Data Moves Through SENTINEL-Q")
add_subtitle(s, "Two streams in. One score out.")
add_bullets(s, [
    "Ingest: security telemetry + transaction logs, identifiers pseudonymized at entry.",
    "Correlate: 5 transparent rules match cyber events to transaction behaviour.",
    "Weight: IsolationForest scores anomalies the rules did not name.",
    "Explain: 0-100 risk score with plain-English reason delivered to the analyst.",
], top=Inches(2.1), height=Inches(2.7))

diagram_box = add_card(s, Inches(1.5), Inches(5.2), Inches(10.3), Inches(1.3))
diagram_box.fill.solid()
diagram_box.fill.fore_color.rgb = TEAL
tf = diagram_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "[INSERT ARCHITECTURE DIAGRAM]"
set_font(run, size=18, bold=True, color=WHITE)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# ---------- SLIDE 15 ----------
s = add_slide()
add_title(s, "See It Live")
add_subtitle(s, "Pick a scenario. Watch it score.")

box_top = Inches(2.3)
box_h = Inches(1.4)
box_w = Inches(3.7)
gap = Inches(0.25)
labels = [("LIVE URL", "sentinel-q-zjrt44jq29pa4uruuhbrkv.streamlit.app"),
          ("DEMO VIDEO", "[VIDEO LINK]"),
          ("REPO", "github.com/Fablescript22/sentinel-q")]
left = Inches(0.8)
for label, content in labels:
    card = add_card(s, left, box_top, box_w, box_h)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    set_font(run, size=14, bold=True, color=TEAL)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = content
    set_font(run2, size=13, bold=False, color=BODY_COLOR)
    left = Emu(left + box_w + gap)

shot_top = Inches(4.1)
shot_h = Inches(2.6)
shot_w = Inches(5.6)
add_placeholder_box(s, Inches(0.8), shot_top, shot_w, shot_h, "[DASHBOARD SCREENSHOT]")
add_placeholder_box(s, Inches(6.7), shot_top, shot_w, shot_h, "[QUANTUM TAB SCREENSHOT]")

# ---------- SLIDE 16 ----------
s = add_slide()
add_title(s, "Built in 48 Hours. Designed to Last.")
add_subtitle(s, "Every claim in this deck runs in the repo.")
add_bullets(s, [
    "One engine that closes the gap between the SOC and the fraud desk.",
    "Transparent rules, explainable scores, honest quantum posture.",
    "Pilot branch to SOC integration to bank-wide. The 30/60/90 starts when a bank says yes.",
])
add_footer(s, "github.com/Fablescript22/sentinel-q | Team: Shreya, Devayani, Meghna, Hitha")

prs.save("docs/SENTINEL_Q_Deck.pptx")
print(f"Saved deck with {len(prs.slides.slides._sldIdLst)} slides" if False else f"Saved deck with {len(prs.slides._sldIdLst)} slides")
